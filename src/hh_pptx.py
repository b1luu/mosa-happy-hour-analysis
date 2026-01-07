from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def get_project_root() -> Path:
    return Path(__file__).resolve().parents[1]


FONT_NAME = "Calibri"
COLOR_PRIMARY = RGBColor(20, 45, 90)
COLOR_ACCENT = RGBColor(30, 130, 76)
COLOR_MUTED = RGBColor(90, 90, 90)


def _style_title(text_frame, size: int = 40) -> None:
    p = text_frame.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = COLOR_PRIMARY


def _style_body(text_frame, size: int = 18) -> None:
    for p in text_frame.paragraphs:
        for run in p.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(size)
            run.font.color.rgb = COLOR_MUTED


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Mosa Tea Happy Hour Performance"
    _style_title(slide.shapes.title.text_frame, size=42)
    subtitle = slide.placeholders[1]
    subtitle.text = "Monday Flat $4 vs Wednesday BOGO 50%"
    _style_body(subtitle.text_frame, size=20)


def _add_text_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    _style_title(slide.shapes.title.text_frame, size=30)
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for text in bullets:
        p = body.add_paragraph()
        p.text = text
        p.level = 0
    _style_body(body, size=18)


def _add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    _style_title(slide.shapes.title.text_frame, size=28)
    left = Inches(0.6)
    top = Inches(1.2)
    height = Inches(5.5)
    slide.shapes.add_picture(str(image_path), left, top, height=height)
    if caption:
        tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(11.0), Inches(0.6))
        tf = tx_box.text_frame
        tf.text = caption
        _style_body(tf, size=14)


def _add_takeaways_slide(prs: Presentation, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Takeaways"
    _style_title(slide.shapes.title.text_frame, size=30)
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for text in bullets:
        p = body.add_paragraph()
        p.text = text
        p.level = 0
    _style_body(body, size=18)


def _add_final_notes_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Final Notes"
    _style_title(slide.shapes.title.text_frame, size=30)
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    notes = [
        "This comparison focuses on net sales and quantity within the Happy Hour window.",
        "BOGO 50% may favor higher volume or pair-order behavior.",
        "Flat $4 may attract price-sensitive or new customers and broaden the audience.",
        "If conclusions feel incomplete, specify the business goal (volume, margin, acquisition, retention).",
    ]
    for text in notes:
        p = body.add_paragraph()
        p.text = text
        p.level = 0
    _style_body(body, size=18)


def _add_section_divider(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12.0), Inches(1.0))
    title_tf = title_box.text_frame
    title_tf.text = title
    _style_title(title_tf, size=36)
    subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.3), Inches(12.0), Inches(0.8))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    _style_body(subtitle_tf, size=18)


def main() -> None:
    project_root = get_project_root()
    plots_dir = project_root / "exports" / "plots"
    output_path = project_root / "exports" / "happy_hour_presentation.pptx"
    summary_by_day = project_root / "exports" / "hh_summary_by_promo_day.csv"
    fruit_summary = project_root / "exports" / "fruit_tea_summary_by_promo_day.csv"
    summary_df = None
    fruit_df = None
    if summary_by_day.exists():
        summary_df = __import__("pandas").read_csv(summary_by_day)
    if fruit_summary.exists():
        fruit_df = __import__("pandas").read_csv(fruit_summary)

    def _get_metric(df, promo, col):
        if df is None or col not in df.columns:
            return None
        row = df[df["promo_day_type"] == promo]
        if row.empty:
            return None
        return float(row.iloc[0][col])

    prs = Presentation()
    _add_title_slide(prs)

    _add_section_divider(prs, "Project Context", "Objectives, scope, and inputs")
    _add_text_slide(
        prs,
        "Project Objective",
        [
            "Compare two Happy Hour promo designs for Fruit Tea: Monday flat $4 vs Wednesday BOGO 50%.",
            "Happy Hour window: Monday & Wednesday, 2:00–5:00 pm (PDT), starting 2025-10-20.",
            "Core question: Which promo drives higher Fruit Tea demand and revenue, and why?",
        ],
    )
    _add_text_slide(
        prs,
        "Data & Pipeline",
        [
            "Source: Square POS item-level exports.",
            "Pipeline: clean text, standardize columns, parse timestamps, add HH flags.",
            "Privacy: remove customer identifiers and sensitive payment fields.",
            "Outputs: processed CSVs, promo summaries, Fruit Tea summaries, product lift.",
        ],
    )
    _add_text_slide(
        prs,
        "Metrics Used",
        [
            "Demand: total Fruit Tea quantity sold during HH window.",
            "Revenue: net sales and effective price per drink (net_total / quantity).",
            "Product mix: top Fruit Tea items per promo day.",
        ],
    )

    _add_section_divider(prs, "Results", "Summary of promo performance")
    image_map = [
        (
            "Happy Hour Net Sales",
            plots_dir / "hh_net_sales_by_day.png",
            "Shows overall revenue for HH window by promo day.",
        ),
        (
            "Happy Hour Total Items",
            plots_dir / "hh_total_items_by_day.png",
            "Shows overall HH volume by promo day (all items).",
        ),
        (
            "Fruit Tea Net Sales",
            plots_dir / "fruit_tea_net_sales_by_day.png",
            "Fruit Tea revenue comparison between the two promos.",
        ),
        (
            "Fruit Tea Quantity",
            plots_dir / "fruit_tea_quantity_by_day.png",
            "Fruit Tea units sold during HH by promo day.",
        ),
        (
            "Fruit Tea Effective Price",
            plots_dir / "fruit_tea_effective_price_by_day.png",
            "Average price paid per Fruit Tea during HH.",
        ),
        (
            "Top Fruit Tea Items: Monday",
            plots_dir / "top_fruit_tea_items_monday_flat_4.png",
            "Top 5 Fruit Tea items on Monday HH.",
        ),
        (
            "Top Fruit Tea Items: Wednesday",
            plots_dir / "top_fruit_tea_items_wednesday_bogo_50.png",
            "Top 5 Fruit Tea items on Wednesday HH.",
        ),
    ]

    for title, path, caption in image_map:
        if path.exists():
            _add_image_slide(prs, title, path, caption)

    _add_section_divider(prs, "Interpretation", "What the results suggest")
    _add_text_slide(
        prs,
        "Interpretation & Strategy",
        [
            "Wednesday BOGO shows higher Fruit Tea quantity and net sales in this period.",
            "Monday flat $4 drives lower effective price per drink (value perception).",
            "Consider BOGO for volume growth; use flat $4 to protect margins or drive value segments.",
            "Promote top Fruit Tea SKUs identified in the product lift charts.",
        ],
    )

    if fruit_df is not None:
        mon_qty = _get_metric(fruit_df, "monday_flat_4", "fruit_tea_quantity")
        wed_qty = _get_metric(fruit_df, "wednesday_bogo_50", "fruit_tea_quantity")
        mon_sales = _get_metric(fruit_df, "monday_flat_4", "fruit_tea_net_sales")
        wed_sales = _get_metric(fruit_df, "wednesday_bogo_50", "fruit_tea_net_sales")
        mon_price = _get_metric(fruit_df, "monday_flat_4", "avg_price_per_item")
        wed_price = _get_metric(fruit_df, "wednesday_bogo_50", "avg_price_per_item")
        bullets = []
        if mon_qty is not None and wed_qty is not None:
            bullets.append(
                f"Fruit Tea volume: Wed {int(round(wed_qty))} vs Mon {int(round(mon_qty))} cups."
            )
        if mon_sales is not None and wed_sales is not None:
            bullets.append(
                f"Fruit Tea net sales: Wed ${wed_sales:,.2f} vs Mon ${mon_sales:,.2f}."
            )
        if mon_price is not None and wed_price is not None:
            bullets.append(
                f"Effective price per drink: Wed ${wed_price:,.2f} vs Mon ${mon_price:,.2f}."
            )
        bullets.append("Recommendation: favor Wednesday BOGO for volume; test Monday flat $4 for value-driven buyers.")
        _add_takeaways_slide(prs, bullets)
    else:
        _add_takeaways_slide(
            prs,
            [
                "Fruit Tea demand and revenue are higher on Wednesday BOGO in this period.",
                "Flat $4 Monday positions value pricing and may appeal to price-sensitive guests.",
                "Use top items to focus signage and staff recommendations.",
            ],
        )

    _add_section_divider(prs, "Notes", "Limitations and next steps")
    _add_final_notes_slide(prs)

    prs.save(output_path)
    print(f"Saved PPTX: {output_path}")


if __name__ == "__main__":
    main()
